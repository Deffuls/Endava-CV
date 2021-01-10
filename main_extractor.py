#GUI Library (Pyqt5)
from PyQt5 import QtWidgets
from PyQt5.QtWidgets import *
from PyQt5.QtGui import *
from PyQt5.QtCore import *
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx import Presentation
from PIL import Image
#Office365 Library for authentication to SHAREPOINT, and for uploading/downloading files easily.
from office365.runtime.auth.authentication_context import AuthenticationContext
from office365.runtime.client_request import ClientRequest
from office365.runtime.utilities.request_options import RequestOptions
from office365.sharepoint.client_context import ClientContext

from office365.runtime.action_type import ActionType
from office365.runtime.client_object import ClientObject
from office365.runtime.client_query import ClientQuery
from office365.runtime.client_result import ClientResult
from office365.runtime.odata.odata_path_parser import ODataPathParser
from office365.runtime.resource_path_entity import ResourcePathEntity
from office365.runtime.resource_path_service_operation import ResourcePathServiceOperation
from office365.runtime.utilities.http_method import HttpMethod
from office365.runtime.utilities.request_options import RequestOptions
from office365.sharepoint.listitem import ListItem
from office365.sharepoint.webparts.limited_webpart_manager import LimitedWebPartManager

import json #for requests
import threading #for threading
import sys
import os
import xlrd #excel library
import pptx #powerpoint library
import requests
import shutil
import datetime
import time

#This is main class "ProcessInfo", this class is processing all information from excel, and
#returning ready string as in "WORD"

class ProcessInfo():

    def __init__(self,project):
        
        self.ProjectName = project
        source = os.getcwd()+"\\"+"resource"+"\\"+"Source.xlsx" #path to Source.xlsx
        self.wb = xlrd.open_workbook(source) #opening excel file
        self.sheet = self.wb.sheet_by_index(0) #getting first page "Project" page.
        
        
    def title(self):
        return self.ProjectName

    def Temperature(self): #This function returns Temperature, from self.ProjectName that we want
        c = self.sheet.nrows
        for x in range(1,c):
            cmp = str(self.sheet.cell_value(x,3))
            if cmp == self.ProjectName: #as you  can see, here it is comparing 'cmp' with 'self.ProjectName'
                s = str(self.sheet.cell_value(x,13))
                return s
                
    
    def RevenueImpact(self): #All functions and this, is returning a excel cell.
        c = self.sheet.nrows #Name of function will tell you what cell will be return
        for x in range(1,c):
            cmp = str(self.sheet.cell_value(x,3))
            if cmp == self.ProjectName:
                
                s = "$x M (Current), $y M (Future), abc (Soft)" #Here we based on information from excel cells
                _x = self.sheet.cell_value(x,7)                 #we are manipulating with string 's'
                _y = self.sheet.cell_value(x,8)
                _abc = self.sheet.cell_value(x,9)
                if not _y:
                    s = s.replace(", $y M (Future)","")
                else:
                    s =s.replace("y",str(_y))
                if not _abc:
                  s = s.replace(", abc (Soft)","")
                else:
                    s = s.replace("abc",str(_abc))

                s = s.replace("x",str(_x))

                return s
#All the functions below are the same at main structure, only difference is that it is returning different string
            #So where is no sense to comment all of them from this class
    def IssueReason(self):
        c = self.sheet.nrows
        for x in range(1,c):
            cmp = str(self.sheet.cell_value(x,3))
            if cmp == self.ProjectName:
                s = str(self.sheet.cell_value(x,36))
                if not s:
                    s = "TBC"

                return s


    def Criteria(self):
        c = self.sheet.nrows
        for x in range(1,c):
            cmp = str(self.sheet.cell_value(x,3))
            if cmp == self.ProjectName:
                s = str(self.sheet.cell_value(x,23))
                if not s:
                    s = "TBC"

                return s

    def RootCause(self):
        c = self.sheet.nrows
        for x in range(1,c):
            cmp = str(self.sheet.cell_value(x,3))
            if cmp == self.ProjectName:
                s = str(self.sheet.cell_value(x,37))
                if not s:
                    s = "TBC"

                return s

    def Lessons(self):
        c = self.sheet.nrows
        for x in range(1,c):
            cmp = str(self.sheet.cell_value(x,3))
            if cmp == self.ProjectName:
                s = str(self.sheet.cell_value(x,38))
                if not s:
                    s = "TBC"

                return s

    def Product(self):
        c = self.sheet.nrows
        for x in range(1,c):
            cmp = str(self.sheet.cell_value(x,3))
            if cmp == self.ProjectName:
                s = str(self.sheet.cell_value(x,15))
                if not s:
                    s = "TBC"

                return s

    def Sponsor(self):
        c = self.sheet.nrows
        for x in range(1,c):
            cmp = str(self.sheet.cell_value(x,3))
            if cmp == self.ProjectName:
                s = str(self.sheet.cell_value(x,17))
                if not s:
                    s = "TBC"

                return s

    def capLeader(self):
        c = self.sheet.nrows
        for x in range(1,c):
            cmp = str(self.sheet.cell_value(x,3))
            if cmp == self.ProjectName:
                s = str(self.sheet.cell_value(x,16))
                if not s:
                    s = "TBC"

                return s

    def AccTeam(self):
        c = self.sheet.nrows
        for x in range(1,c):
            cmp = str(self.sheet.cell_value(x,3))
            if cmp == self.ProjectName:
                s = str(self.sheet.cell_value(x,18))
                if not s:
                    s = "TBC"

                return s

    def CoreTeam(self):
        self.sheet = self.wb.sheet_by_index(1) #changing to tab "Actions"
        c = self.sheet.nrows
        owners = []
        for x in range(0,c):
            cmp = str(self.sheet.cell_value(x,0))
            if cmp == self.ProjectName:
                apnd = str(self.sheet.cell_value(x,5))
                bad_chars = [';','#','0','1','2','3'  #here we are cheching for bad characters and removing them
                             ,'4','5','6','7','8','9']
                for i in bad_chars: #looping 'i' - is equal to first element of array, and checking in function if we have one
                    apnd = apnd.replace(i,'') #here we are removing if we found one
                                
                if not apnd in owners:
                    owners.append(apnd)
                        
        self.sheet = self.wb.sheet_by_index(0) #changing back to tab "Project"
        if not owners:
            return "TBC"
        else:
            return owners #here we are returning array of names

    def Customers(self):
        c = self.sheet.nrows
        owners = []
        for x in range(1,c):
            cmp = str(self.sheet.cell_value(x,3))
            if cmp == self.ProjectName:
                s = str(self.sheet.cell_value(x,20))
                if not s:
                    s = "TBC"

                return s

    def StartDate(self):
        c = self.sheet.nrows
        for x in range(1,c):
            cmp = str(self.sheet.cell_value(x,3))
            if cmp == self.ProjectName:
                s = datetime.datetime(*xlrd.xldate_as_tuple(self.sheet.cell_value(x,10),self.wb.datemode)) #getting from excel cell date, and converting
                return str(s)[:10] # here returning only DATE XX/ZZ/DDDD, without TIME, in second,minutes,hours
                                    #it to a string by a special function from xlrd

    def CloseDate(self): #this function is the same as StartDate
        c = self.sheet.nrows
        for x in range(1,c):
            cmp = str(self.sheet.cell_value(x,3))
            if cmp == self.ProjectName:
                s = datetime.datetime(*xlrd.xldate_as_tuple(self.sheet.cell_value(x,11),self.wb.datemode))
                return str(s)[:10]

    def DaysOnCap(self,startdate,closedate): #take 2 parameters startdate and closedate, and checking what to return as last string
        if not closedate:
            closedate = startdate
            return (closedate + " - " + startdate)
        else:
            now = datetime.date.today()
            return (str(now) + " - " + startdate)

    def BauActions(self):
        c = self.sheet.nrows
        for x in range(1,c):
            cmp = str(self.sheet.cell_value(x,3))
            if cmp == self.ProjectName:
                s = str(self.sheet.cell_value(x,39))
                if not s:
                    s = "TBC"

                return s

    def GetUrl(self):
        c = self.sheet.nrows
        for x in range(1,c):
            cmp = str(self.sheet.cell_value(x,3))
            if cmp == self.ProjectName:
                s = str(self.sheet.cell_value(x,35))
                if s == None:
                    s = None
                
                return s







    

def Folders(): #this functions checks if we have all folders necessarry for doing work, if it don't exist it will create one
    path = os.getcwd()
    if os.path.isdir(path+"\\"+"resource") == False: #resources we have Source.xlsx and Template.pptx
        os.mkdir(path+"\\"+"resource")
    if os.path.isdir(path+"\\"+"logos") == False: #and in logos we can found logos from every .pptx file.
        os.mkdir(path+"\\"+"logos")       
    
        

def GetFiles(login,password,url_login):
    ctx_auth = AuthenticationContext(url_login) #creating ctx_auth session with function AuthenticationContext(by a url_login), look in settings.txt
    ctx_auth.acquire_token_for_user(login, password) #getting token auth by loggin into site
    ctx = ClientContext(url_login, ctx_auth)#and finally logging in
    filename="Template.pptx"
    response = open_binary(ctx, "/Shared Documents/"+filename,filename) #getting ctx connection to a location on sharepoint to download files
        
    with open(os.getcwd()+"\\"+"resource\\"+filename, "wb") as local_file: #writing files from response, in resource folder.in BINARY MODE.
        local_file.write(response.content) #response.content, it is binary information of files, Template.pptx

        
    filename="Source.xlsx"
    response = open_binary(ctx, "/Shared Documents/"+filename,filename) #getting ctx connection to a location on sharepoint to download files
    with open(os.getcwd()+"\\"+"resource\\"+filename, "wb") as local_file: #writing files from response, in resource folder.in BINARY MODE.
        local_file.write(response.content) #response.content, it is binary information of files, Source.xlsx

def InitializeComponents(login,password,url_login):
    t1 = threading.Thread(target=Folders)
    t1.start()
    t2 = threading.Thread(target=GetFiles,args=(login,password,url_login,))
    t1.join()
        
    t2.start()
        
    t2.join()
    print("Initialize finished.")
    

def UploadFile(filename,path,login,password,url_login):
    print("UPLOAD FILE --------")
    EMAIL = login #set email(login)
    PASSWORD = password# set password
    url_login = url_login #set login
        
        
    ctx_auth = AuthenticationContext(url_login) #authentication by url
    ctx_auth.acquire_token_for_user(EMAIL, PASSWORD)    #creating token 
    ctx = ClientContext(url_login, ctx_auth) #initialize login finally

        
        
    with open(path+"\\"+filename, 'rb') as content_file: #reading file to upload it to sharepoint
        file_content = content_file.read()
        #parameters:  ctx, shared documents is folder,file_content in binary, and filename "Example.pptx"
        r = save_binary(ctx,'Shared Documents/',file_content,filename) # and saving it on sharepoint 
        

def GeneratePPTX(homeDir,info,login,password,url_login): #Generating pptx function
    print("GENERATING PPTX!")
    path = homeDir #getting homedir
    print("PATH HOMEDIR: "+ path)
     #getting title of .pptx, and name to find information in excel
    print("2: " + path+"\\"+"resource\\Template.pptx")
    prs = Presentation(path+"\\"+"resource\\Template.pptx") #openning Template to generate .pptx
                
                
    slide_layout = prs.slide_layouts[7] #set layout, template, from position 7

    slide = prs.slides.add_slide(slide_layout) #creating a slide with layout 7
    shapes = slide.shapes #get all shapes from this slide
                
    shapes[0].text = info.ProjectName #setting Project name by Shape, no matter how, we can do it also like a placeholder, just less code
    url = info.GetUrl() #checking for url to download LOGO
    name_f = None
    if not url == None:
        name_f = GetFileName(url) #call function that gets file name of logo
    if not name_f == None: #checking if exist or not, if exists executing code below, if not ignore it
        if not str(name_f[len(name_f)-8:len(name_f)-4]) == ".svg": #
            try:
                DownloadFile(path+"\\"+"logos"+"\\"+name_f,url) #downloading file to logos folder
                img = Image.open(path+"\\"+"logos"+"\\"+name_f) #openiing it
                            
                width, height = img.size #resizing it to the height of shape
                slide.shapes[1].height = height
                slide.shapes[1].width = width
                            
                placeholder = slide.shapes[1].insert_picture(path+"\\"+"logos"+"\\"+name_f) # inserting it to Picture placeholder
                            
                image_ratio = width / height
                placeholder_ratio = placeholder.width / placeholder.height
                ratio_difference = placeholder_ratio - image_ratio
                             #some calculations for having right proportions of image in placeholder
                            # Placeholder width too wide:
                if ratio_difference > 0:
                    difference_on_each_side = ratio_difference / 2
                    placeholder.crop_left = -difference_on_each_side
                    placeholder.crop_right = -difference_on_each_side
                            # Placeholder height too high
                else:
                    difference_on_each_side = -ratio_difference / 2
                    placeholder.crop_bottom = -difference_on_each_side
                    placeholder.crop_top = -difference_on_each_side
                    
                            
            except requests.exceptions.RequestException as e:
                exit(0)

    temp = info.Temperature() #Getting temperature
    if temp == "Cool": 
        temperature = shapes[17]
        temperature.fill.solid() #set shape solid
        temperature.fill.fore_color.rgb = RGBColor(0, 128, 0) #and color it to respective color
    elif temp == "Hot":
        temperature = shapes[17]
        temperature.fill.solid()
        temperature.fill.fore_color.rgb = RGBColor(255, 0, 0)
    elif temp == "At Risk":
        temperature = shapes[17]
        temperature.fill.solid()
        temperature.fill.fore_color.rgb = RGBColor(255, 255, 0)
                #MAIN TABLE -------------------------------------------

    add_pptx(shapes[5],10,info.RevenueImpact(),1) #this function adding text, to respective placeholder
    add_pptx(shapes[6],10,info.Product(),1)
    add_pptx(shapes[2],10,info.IssueReason(),1)
    add_pptx(shapes[3],10,info.Criteria(),1)
    add_pptx(shapes[7],10,info.RootCause(),1)
    add_pptx(shapes[8],10,info.Lessons(),1)
    add_pptx(shapes[4],10,info.BauActions(),1)

                #MAIN TABLE -------------------------------------------
                #SECOND TABLE -------------------------------------------
    add_pptx(shapes[9],9,info.Sponsor(),1)
    add_pptx(shapes[10],9,info.capLeader(),1)
    add_pptx(shapes[11],9,info.AccTeam(),1)
    add_pptx(shapes[12],9,info.CoreTeam(),0)
    add_pptx(shapes[13],9,info.Customers(),1)
    add_pptx(shapes[14],9,info.StartDate(),1)
    add_pptx(shapes[15],9,info.CloseDate(),1)
    add_pptx(shapes[16],9,info.DaysOnCap(info.StartDate(),info.CloseDate()),1)

    path = os.path.join(os.path.join(os.environ['USERPROFILE']), 'Desktop') #getting desktop location
    print("path desktop: " + path)
    print("path saved: "+ path + "\\"+ info.ProjectName + ".pptx")
    prs.save(path + "\\"+ info.ProjectName + ".pptx") #saving .pptx file to desktop by name
    print("save: " + path + "\\"+ info.ProjectName + ".pptx")
    print("string to upload: " + info.ProjectName + ".pptx",path)
    UploadFile(info.ProjectName + ".pptx",path,login,password,url_login) #and uploading it to sharepoint   def UploadFile(filename,path,login,password,url_login):
    path = homeDir #setting path to homeDir to avoid problems in future

            
def open_binary(ctx, server_relative_url,filename):
    print("OPEN BINARY ----------")
    try:
        from urllib import quote  
    except ImportError:
        from urllib.parse import quote  
    server_relative_url = quote(server_relative_url,filename) #sets relativeurl
    url = ctx.service_root_url+"web/GetFolderByServerRelativeUrl('Shared Documents/')/Files('"+filename+"')/$value" #main url where is file
    request = RequestOptions(url)
    request.method = HttpMethod.Get
    response = ctx.execute_request_direct(request)
    print(response)
    return response
    
def save_binary(ctx, server_relative_url, content,filename):
    try:
        from urllib import quote  # Python 2.X
    except ImportError:
        from urllib.parse import quote  # Python 3+

    server_relative_url = quote(server_relative_url)
    url = ctx.service_root_url+"web/GetFolderByServerRelativeUrl('Shared Documents/')/Files/add(url='"+filename+"',overwrite=true)" #main url where to save file

    request = RequestOptions(url) #setting request url
    request.method = HttpMethod.Post #method post
    request.set_header('X-HTTP-Method-Override', 'PUT') #set header Ovveride and Put
    request.data = content #content, BINARY DATA OF FILE
    response = ctx.execute_request_direct(request) #and executing it 
    print(response)
    return response



def add_pptx(par,size,string,menu):
    fill = par.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255,255,255)
    fill.background()
    if menu == 1: # two menu, if 1 then it is only string 
        par.text = string
    elif menu == 0: # if 0 then it is array (list)
        for x in range(0,len(string)): 
            par.text += string[x]

def GetFileName(page):
    if page == None:
        return None
    c = len(page)
    if not c == 0 or c < 0:
        c -= 1
        for x in range(c,0,-1): #start loop from end of string to start of string
            if page[x] == '/': #find first slash '/', and from index of this slash we are taking the name of file, example.png,jpeg,...
                return page[x+1:]
    else:
        return None

def DownloadFile(name,url): #download function for logos,#creating file with name of logo, and writing it in binary mode
    with open(name,'wb') as file,\
     requests.get(url, stream=True) as response: #requesting file by get request, adn streaming data as response
     shutil.copyfileobj(response.raw,file) #writing it to a file 
            

        

if __name__ == '__main__':
    with open('settings.txt', 'r') as f:
        login = f.readline().replace('\n', '') #deleting '\n', to have the string in one line not from new
        password = f.readline().replace('\n', '')
        url_login = f.readline().replace('\n', '')

        login = str(login.replace('login:', '')) #deleting from each string what we don't need
        password = str(password.replace('password:', ''))
        url_login = str(url_login.replace('share_url:', ''))
    InitThread = threading.Thread(target=InitializeComponents,args=(login,password,url_login,)) #Starting a thread, for function InitalizeComponents
    print("Initialzie components")
    InitThread.start()
    InitThread.join()
    ProjectName = input('Project name: ')
    ProjectName = str(ProjectName)
    info = ProcessInfo(ProjectName)
    pptx_thread = threading.Thread(target=GeneratePPTX,args=(os.getcwd(),info,login,password,url_login)) #by pressing GenerateButton call this function
    pptx_thread.start() #starting thread
    pptx_thread.join() #waiting it to join



   #Arab National Bank (ANB)
   #Axis Bank - CP-eM CAP
   #Royal Bank of Canada
   #Banc First
   #Nedbank
   #
   #Trader Joe’s

